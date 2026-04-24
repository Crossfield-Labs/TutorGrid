from __future__ import annotations

from pathlib import Path
from tempfile import TemporaryDirectory

from backend.knowledge.parsers.base import ParsedBlock, ParsedDocument
from backend.knowledge.parsers.image_ocr_parser import ImageOcrParser


class PptxParser:
    def __init__(self, *, enable_image_ocr: bool = True) -> None:
        self.enable_image_ocr = enable_image_ocr
        self.ocr_parser = ImageOcrParser() if enable_image_ocr else None

    def parse(self, file_path: Path) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover
            raise RuntimeError("python-pptx is required for PPTX parsing.") from exc

        presentation = Presentation(str(file_path))
        blocks: list[ParsedBlock] = []
        with TemporaryDirectory(prefix="orchestrator_pptx_ocr_") as temp_dir:
            temp_root = Path(temp_dir)
            for slide_index, slide in enumerate(presentation.slides, start=1):
                texts: list[str] = []
                image_blocks: list[str] = []
                for shape_index, shape in enumerate(slide.shapes, start=1):
                    if getattr(shape, "has_text_frame", False):
                        text = str(shape.text or "").strip()
                        if text:
                            texts.append(text)
                    ocr_text = self._extract_shape_image_text(
                        shape=shape,
                        temp_root=temp_root,
                        slide_index=slide_index,
                        shape_index=shape_index,
                    )
                    if ocr_text:
                        image_blocks.append(ocr_text)

                merged = "\n".join(texts).strip()
                if merged:
                    blocks.append(
                        ParsedBlock(
                            text=merged,
                            section=f"slide-{slide_index}",
                            page=slide_index,
                            metadata={"slideIndex": slide_index},
                        )
                    )
                for image_index, image_text in enumerate(image_blocks, start=1):
                    blocks.append(
                        ParsedBlock(
                            text=image_text,
                            section=f"slide-{slide_index}-image-{image_index}",
                            page=slide_index,
                            metadata={
                                "slideIndex": slide_index,
                                "imageIndex": image_index,
                                "source": "image_ocr",
                            },
                        )
                    )
        return ParsedDocument(
            title=file_path.stem,
            blocks=blocks,
            metadata={"parser": "python-pptx", "sourcePath": str(file_path)},
        )

    def _extract_shape_image_text(
        self,
        *,
        shape: object,
        temp_root: Path,
        slide_index: int,
        shape_index: int,
    ) -> str:
        if self.ocr_parser is None:
            return ""
        image = getattr(shape, "image", None)
        if image is None:
            return ""
        blob = getattr(image, "blob", b"")
        if not blob:
            return ""
        ext = str(getattr(image, "ext", "png") or "png").strip(".").lower() or "png"
        file_path = temp_root / f"slide_{slide_index}_shape_{shape_index}.{ext}"
        try:
            file_path.write_bytes(blob)
            lines = self.ocr_parser.extract_lines(file_path)
        except Exception:
            return ""
        return "\n".join(lines).strip()
